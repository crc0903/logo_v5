import streamlit as st
from PIL import Image, ImageChops
import os
import io
import math
from pptx import Presentation
from pptx.util import Inches

PRELOADED_LOGO_DIR = "preloaded_logos"

def load_preloaded_logos():
    logos = {}
    if not os.path.exists(PRELOADED_LOGO_DIR):
        os.makedirs(PRELOADED_LOGO_DIR)
    for file in os.listdir(PRELOADED_LOGO_DIR):
        if file.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
            name = os.path.splitext(file)[0]
            image = Image.open(os.path.join(PRELOADED_LOGO_DIR, file)).convert("RGBA")
            logos[name] = image
    return logos

# Trim white or transparent space around the logo
def trim_whitespace(image):
    bg = Image.new(image.mode, image.size, (255, 255, 255, 0))  # transparent background
    diff = ImageChops.difference(image, bg)
    bbox = diff.getbbox()
    if bbox:
        return image.crop(bbox)
    return image

# Resize logo to fill a dynamic 5x2 box inside the cell (whichever is limiting)
def resize_to_fill_5x2_box(image, cell_width_px, cell_height_px, buffer_ratio=0.9):
    box_ratio = 5 / 2
    max_box_width = int(cell_width_px * buffer_ratio)
    max_box_height = int(cell_height_px * buffer_ratio)

    # Fit the largest possible 5x2 box inside the cell
    if max_box_width / box_ratio <= max_box_height:
        box_width = max_box_width
        box_height = int(max_box_width / box_ratio)
    else:
        box_height = max_box_height
        box_width = int(max_box_height * box_ratio)

    # Resize logo proportionally to fit within that box
    img_w, img_h = image.size
    img_ratio = img_w / img_h

    if img_ratio > (box_width / box_height):
        new_width = box_width
        new_height = int(box_width / img_ratio)
    else:
        new_height = box_height
        new_width = int(box_height * img_ratio)

    resized = image.resize((new_width, new_height), Image.LANCZOS)
    return resized, box_width, box_height

def create_logo_slide(prs, logos, canvas_width_in, canvas_height_in, logos_per_row):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas_width_px = int(canvas_width_in * 96)
    canvas_height_px = int(canvas_height_in * 96)

    logo_count = len(logos)
    cols = logos_per_row if logos_per_row else max(1, round(math.sqrt(logo_count * canvas_width_in / canvas_height_in)))
    rows = math.ceil(logo_count / cols)

    cell_width = canvas_width_px / cols
    cell_height = canvas_height_px / rows

    left_margin = Inches((10 - canvas_width_in) / 2)
    top_margin = Inches((7.5 - canvas_height_in) / 2)

    for idx, logo in enumerate(logos):
        col = idx % cols
        row = idx // cols

        trimmed = trim_whitespace(logo)
        resized, box_width, box_height = resize_to_fill_5x2_box(trimmed, int(cell_width), int(cell_height))

        img_stream = io.BytesIO()
        resized.save(img_stream, format="PNG")
        img_stream.seek(0)

        # Center the logo inside the 5x2 box, and box inside the cell
        x_offset = (cell_width - box_width) / 2 + (box_width - resized.width) / 2
        y_offset = (cell_height - box_height) / 2 + (box_height - resized.height) / 2
        left = left_margin + Inches((col * cell_width + x_offset) / 96)
        top = top_margin + Inches((row * cell_height + y_offset) / 96)

        slide.shapes.add_picture(
            img_stream, left, top,
            width=Inches(resized.width / 96),
            height=Inches(resized.height / 96)
        )

# --- Streamlit UI ---
st.title("Logo Grid PowerPoint Exporter")
st.markdown("Upload logos or use preloaded ones below:")

uploaded_files = st.file_uploader("Upload logos", type=["png", "jpg", "jpeg", "webp"], accept_multiple_files=True)
preloaded = load_preloaded_logos()
selected_preloaded = st.multiselect("Select preloaded logos", options=list(preloaded.keys()))

canvas_width_in = st.number_input("Grid width (inches)", min_value=1.0, max_value=20.0, value=10.0)
canvas_height_in = st.number_input("Grid height (inches)", min_value=1.0, max_value=20.0, value=7.5)
logos_per_row = st.number_input("Logos per row (optional)", min_value=0, max_value=50, value=0)

if st.button("Generate PowerPoint"):
    images = []

    if uploaded_files:
        for f in uploaded_files:
            image = Image.open(f).convert("RGBA")
            images.append(image)

    for name in sorted(selected_preloaded):
        images.append(preloaded[name])

    if not images:
        st.warning("Please upload or select logos.")
    else:
        prs = Presentation()
        create_logo_slide(prs, images, canvas_width_in, canvas_height_in,
                          logos_per_row if logos_per_row > 0 else None)

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        st.success("PowerPoint created!")
        st.download_button("Download .pptx", output, file_name="logo_grid.pptx")
